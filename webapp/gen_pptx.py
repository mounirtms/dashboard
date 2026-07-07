#!/usr/bin/env python3.9
"""
TechnoStationery Executive Audit — PPTX Generator
12-slide condensed version from the 25-slide HTML presentation
Professional dark theme matching the dashboard branding
"""
import sys, os
sys.path.insert(0, '/home/dashboard/public_html/webapp/pptx_lib')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import io

# ── Palette ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0a, 0x0e, 0x1a)   # main dark bg
BG2      = RGBColor(0x0f, 0x17, 0x2a)   # panel bg
ACCENT   = RGBColor(0x3b, 0x82, 0xf6)   # blue
ACCENT2  = RGBColor(0x06, 0xb6, 0xd4)   # cyan
GREEN    = RGBColor(0x22, 0xc5, 0x5e)
YELLOW   = RGBColor(0xf5, 0x9e, 0x0b)
RED      = RGBColor(0xef, 0x44, 0x44)
PURPLE   = RGBColor(0x8b, 0x5c, 0xf6)
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
DIM      = RGBColor(0x47, 0x55, 0x69)
PANEL    = RGBColor(0x11, 0x18, 0x27)
BORDER   = RGBColor(0x1e, 0x2d, 0x45)

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank

def add_slide(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    # Full-bleed dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def txb(slide, text, x, y, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a text box and return the shape."""
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf

def rect(slide, x, y, w, h, fill_color=PANEL, line_color=BORDER, line_w=Pt(0.75)):
    """Add a filled rounded rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def hline(slide, x, y, w, color=BORDER):
    """Thin horizontal rule."""
    line = slide.shapes.add_shape(1, x, y, w, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def section_badge(slide, text, x, y, color=ACCENT):
    """Small uppercase section label."""
    r = rect(slide, x, y, Inches(1.8), Inches(0.22),
             fill_color=RGBColor(
                 int(color[0]*0.15), int(color[1]*0.15), int(color[2]*0.15)
             ) if hasattr(color, '__getitem__') else PANEL,
             line_color=color, line_w=Pt(0.5))
    txb(slide, text, x + Inches(0.05), y, Inches(1.7), Inches(0.22),
        size=7, bold=True, color=color, align=PP_ALIGN.CENTER)
    return r

def kpi_box(slide, label, value, unit, x, y, w=Inches(2.2), h=Inches(1.0), color=ACCENT):
    rect(slide, x, y, w, h, fill_color=RGBColor(0x0d, 0x1a, 0x2e), line_color=color, line_w=Pt(0.75))
    txb(slide, value, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), Inches(0.5),
        size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, label, x + Inches(0.1), y + Inches(0.58), w - Inches(0.2), Inches(0.22),
        size=8, bold=False, color=MUTED, align=PP_ALIGN.CENTER)
    if unit:
        txb(slide, unit, x + Inches(0.1), y + Inches(0.78), w - Inches(0.2), Inches(0.18),
            size=7, bold=False, color=DIM, align=PP_ALIGN.CENTER)

def row_item(slide, label, value, x, y, w=Inches(4.0), bar_pct=None, bar_color=ACCENT):
    """A labeled metric row with optional bar."""
    txb(slide, label, x, y, w * 0.65, Inches(0.22), size=8.5, color=MUTED)
    txb(slide, value, x + w * 0.65, y, w * 0.35, Inches(0.22), size=8.5,
        bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    if bar_pct is not None:
        # bg bar
        rect(slide, x, y + Inches(0.21), w, Inches(0.06),
             fill_color=BORDER, line_color=None)
        # filled bar
        rect(slide, x, y + Inches(0.21), w * (bar_pct / 100), Inches(0.06),
             fill_color=bar_color, line_color=None)

def bullet_list(slide, items, x, y, w, spacing=0.28, size=9, color=MUTED, dot_color=ACCENT):
    """Render a bulleted list."""
    for i, item in enumerate(items):
        cy = y + Inches(spacing * i)
        # bullet dot
        dot = slide.shapes.add_shape(1, x, cy + Inches(0.07), Inches(0.06), Inches(0.06))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        txb(slide, item, x + Inches(0.12), cy, w - Inches(0.12), Inches(0.25),
            size=size, color=color)

def page_header(slide, section, title, subtitle=None):
    """Standard slide header with section badge, title, subtitle."""
    # Top accent bar
    rect(slide, 0, 0, W, Inches(0.06), fill_color=ACCENT, line_color=None)
    # Section label
    txb(slide, f"  {section}  ", Inches(0.35), Inches(0.15), Inches(3), Inches(0.24),
        size=7.5, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    # Title
    txb(slide, title, Inches(0.35), Inches(0.38), W - Inches(3.5), Inches(0.55),
        size=22, bold=True, color=WHITE)
    if subtitle:
        txb(slide, subtitle, Inches(0.35), Inches(0.93), W - Inches(3.5), Inches(0.28),
            size=9, color=MUTED, italic=False)
    # Page number placeholder (right side)
    txb(slide, "technostationery.com", W - Inches(2.1), Inches(0.18), Inches(2.0), Inches(0.22),
        size=7, color=DIM, align=PP_ALIGN.RIGHT)
    hline(slide, Inches(0.35), Inches(1.2), W - Inches(0.7))

def footer_bar(slide, slide_num, total=12):
    hline(slide, 0, H - Inches(0.3), W, color=BORDER)
    txb(slide, "TechnoStationery Executive Audit · Jan–Jul 2026 · Confidential",
        Inches(0.35), H - Inches(0.28), W * 0.6, Inches(0.22),
        size=6.5, color=DIM)
    txb(slide, f"{slide_num} / {total}",
        W - Inches(0.8), H - Inches(0.28), Inches(0.7), Inches(0.22),
        size=6.5, color=DIM, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x05, 0x08, 0x10)

# Gradient accent shapes (simulated with translucent rects)
r = slide.shapes.add_shape(1, W - Inches(4), Inches(-0.5), Inches(5), Inches(5))
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x0a, 0x1a, 0x3a); r.line.fill.background()

r2 = slide.shapes.add_shape(1, Inches(-1), H - Inches(4), Inches(5), Inches(5))
r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x05, 0x12, 0x20); r2.line.fill.background()

# Top accent line
rect(slide, 0, 0, W, Inches(0.05), fill_color=ACCENT, line_color=None)
rect(slide, 0, 0, Inches(3), Inches(0.05), fill_color=ACCENT2, line_color=None)

# Company label
txb(slide, "TECHNOSTATIONERY.COM", Inches(0.5), Inches(0.8), W - Inches(1), Inches(0.35),
    size=10, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# Main title
txb(slide, "Executive Audit Report", Inches(0.5), Inches(1.35), W - Inches(1), Inches(0.9),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle — date range
txb(slide, "January – July 2026", Inches(0.5), Inches(2.2), W - Inches(1), Inches(0.45),
    size=24, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

hline(slide, Inches(3.5), Inches(2.75), Inches(6.33))

# Tagline
txb(slide, "Infrastructure · Security · Performance · Business Intelligence",
    Inches(0.5), Inches(2.9), W - Inches(1), Inches(0.3),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)
txb(slide, "8-Phase Forensic Methodology · Evidence-First · Cross-Validated",
    Inches(0.5), Inches(3.2), W - Inches(1), Inches(0.3),
    size=10, color=DIM, align=PP_ALIGN.CENTER)

# KPI strip
kpis = [
    ("Uptime",        "99.7%",  "6 months",  GREEN),
    ("Orders",        "1,247",  "Jan–Jun 26", ACCENT),
    ("Security",      "A+",     "Rating",     ACCENT2),
    ("Commits",       "96",     "H1 2026",    PURPLE),
    ("Bugs Fixed",    "32",     "Sprint",     YELLOW),
    ("Server Load",   "↓68%",  "vs H1'25",   GREEN),
]
kpi_x = Inches(0.5)
kw = (W - Inches(1)) / len(kpis)
for i, (lbl, val, unit, col) in enumerate(kpis):
    kpi_box(slide, lbl, val, unit, kpi_x + kw * i + Inches(0.05),
            Inches(3.75), kw - Inches(0.1), Inches(1.05), col)

# Presented by
txb(slide, "Prepared by: Mounir Abderrahmani  ·  Full-Stack DevOps · eCommerce Architect",
    Inches(0.5), Inches(5.05), W - Inches(1), Inches(0.28),
    size=8.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
txb(slide, "CONFIDENTIAL — For Internal Executive Use Only",
    Inches(0.5), Inches(5.3), W - Inches(1), Inches(0.28),
    size=8, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Bottom meta strip
rect(slide, 0, H - Inches(0.55), W, Inches(0.55),
     fill_color=RGBColor(0x03, 0x05, 0x0c), line_color=None)
rect(slide, 0, H - Inches(0.55), W, Inches(0.02),
     fill_color=BORDER, line_color=None)
metas = ["Platform: Magento 2.4.7", "Server: AlmaLinux 9.6", "DB: MariaDB 10.6", "Cache: Redis + Varnish + Cloudflare"]
meta_w = W / len(metas)
for i, m in enumerate(metas):
    txb(slide, m, meta_w * i + Inches(0.1), H - Inches(0.42),
        meta_w - Inches(0.1), Inches(0.3), size=8, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY & KPIs
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "01 · EXECUTIVE OVERVIEW", "Key Performance Indicators — Jan–Jul 2026",
            "At-a-glance health scorecard across all monitored dimensions")
footer_bar(slide, 2)

# 8 KPI boxes in 2 rows
kpi_data = [
    ("Server Uptime",       "99.7%",  "6-Month Avg",      GREEN),
    ("Avg Load (1min)",     "1.8",    "Normal: <4.0",      GREEN),
    ("Memory Used",         "42%",    "4GB / 9.7GB",      GREEN),
    ("Orders (Jan–Jun)",    "1,247",  "+12% vs H1'25",    ACCENT),
    ("Revenue Growth",      "+18%",   "YoY 2025→2026",    GREEN),
    ("Cache Hit Ratio",     "87%",    "Varnish+Redis",    ACCENT2),
    ("Security Score",      "A+",     "Post-hardening",   PURPLE),
    ("SSH Attacks Blocked", "14.2K",  "Jun 8–14 peak",    RED),
]
cols = 4
kw2 = Inches(2.9)
kh2 = Inches(0.95)
gx  = Inches(0.35)
gy  = Inches(1.35)
gap_x = (W - gx * 2 - kw2 * cols) / (cols - 1)
for i, (lbl, val, unit, col) in enumerate(kpi_data):
    row = i // cols
    c   = i % cols
    x = gx + c * (kw2 + gap_x)
    y = gy + row * (kh2 + Inches(0.15))
    kpi_box(slide, lbl, val, unit, x, y, kw2, kh2, col)

# Bottom summary bar
rect(slide, Inches(0.35), Inches(5.45), W - Inches(0.7), Inches(1.4),
     fill_color=RGBColor(0x06, 0x10, 0x20), line_color=ACCENT, line_w=Pt(0.5))
txb(slide, "AUDIT SCOPE SUMMARY",
    Inches(0.5), Inches(5.55), W - Inches(1), Inches(0.22),
    size=8, bold=True, color=ACCENT)
scope = [
    "✓  Infrastructure forensics: server hardware, OS config, services, resource utilization",
    "✓  Security audit: SSH brute-force (14,218 attempts), CVE matrix, Imunify360 analysis, hardening before/after",
    "✓  Commerce BI: 1,247 orders, 6-month trend, wilaya breakdown (48 provinces), YoY comparison",
    "✓  Performance deep-dive: Varnish cache, Redis, Cloudflare CDN, query optimization, May crisis post-mortem",
]
for i, s in enumerate(scope):
    txb(slide, s, Inches(0.5), Inches(5.78) + Inches(0.23 * i),
        W - Inches(1), Inches(0.22), size=7.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — DEVELOPMENT & GIT ANALYTICS
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "02 · DEVELOPMENT", "Git Analytics & Sprint KPIs — H1 2026",
            "Commit velocity, feature delivery, and sprint completion rates")
footer_bar(slide, 3)

# Left: commit stats
rect(slide, Inches(0.35), Inches(1.35), Inches(4.2), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "COMMIT ANALYTICS", Inches(0.5), Inches(1.5), Inches(4), Inches(0.25),
    size=8, bold=True, color=ACCENT)
hline(slide, Inches(0.5), Inches(1.78), Inches(3.9), color=BORDER)

git_stats = [
    ("Total Commits (H1)",    "96",    None,  ACCENT),
    ("Bugs Fixed",            "32",    None,  RED),
    ("Features Delivered",    "38",    None,  GREEN),
    ("Tasks Closed",          "20",    None,  YELLOW),
    ("Active Branches",       "3",     None,  PURPLE),
    ("Code Coverage (est.)",  "~72%",  None,  ACCENT2),
]
for i, (lbl, val, _, col) in enumerate(git_stats):
    y = Inches(1.95) + Inches(0.48 * i)
    rect(slide, Inches(0.5), y, Inches(3.8), Inches(0.38),
         fill_color=RGBColor(0x0d, 0x1a, 0x2d), line_color=BORDER, line_w=Pt(0.5))
    txb(slide, lbl, Inches(0.65), y + Inches(0.08), Inches(2.4), Inches(0.25),
        size=9, color=MUTED)
    txb(slide, val, Inches(3.6), y + Inches(0.08), Inches(0.65), Inches(0.25),
        size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)

# Right: sprint progress bars
rect(slide, Inches(4.8), Inches(1.35), Inches(8.18), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "SPRINT PROGRESS", Inches(4.95), Inches(1.5), Inches(5), Inches(0.25),
    size=8, bold=True, color=GREEN)
hline(slide, Inches(4.95), Inches(1.78), Inches(7.9), color=BORDER)

sprints = [
    ("Build & Deploy Pipeline",  100, GREEN),
    ("Logo & Brand Rollout",     100, ACCENT),
    ("Security Hardening",        85, YELLOW),
    ("Performance Tuning",        70, ACCENT2),
    ("Magento Commerce Pages",    80, PURPLE),
    ("Presentation & Reporting", 100, RGBColor(0xec, 0x48, 0x99)),
    ("Cache Optimization",        90, GREEN),
    ("Monitoring & Alerts",       75, ACCENT),
]
for i, (label, pct, col) in enumerate(sprints):
    y = Inches(1.95) + Inches(0.43 * i)
    txb(slide, label, Inches(4.95), y, Inches(5.3), Inches(0.22), size=8.5, color=MUTED)
    txb(slide, f"{pct}%", Inches(10.7), y, Inches(0.55), Inches(0.22),
        size=8.5, bold=True, color=col, align=PP_ALIGN.RIGHT)
    # bg bar
    rect(slide, Inches(4.95), y + Inches(0.22), Inches(6.0), Inches(0.08),
         fill_color=BORDER, line_color=None)
    # progress fill
    if pct > 0:
        rect(slide, Inches(4.95), y + Inches(0.22), Inches(6.0 * pct / 100), Inches(0.08),
             fill_color=col, line_color=None)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — SERVER INFRASTRUCTURE
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "03 · INFRASTRUCTURE", "Server Hardware & OS Configuration",
            "AlmaLinux 9.6 · 8-core Xeon · 9.7GB RAM · NVMe SSD")
footer_bar(slide, 4)

# Left panel - Hardware specs
rect(slide, Inches(0.35), Inches(1.35), Inches(4.0), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "HARDWARE SPECS", Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.25),
    size=8, bold=True, color=ACCENT)
hline(slide, Inches(0.5), Inches(1.78), Inches(3.7))

hw = [
    ("OS",           "AlmaLinux 9.6.20250507"),
    ("Kernel",       "5.14.0-503.el9"),
    ("CPU",          "Intel Xeon @ 2.40GHz × 8"),
    ("RAM",          "9.7GB total / 4.0GB used"),
    ("Disk",         "NVMe SSD 186GB / 59% used"),
    ("Swap",         "4.0GB / minimal usage"),
    ("Network",      "GigE · Cloudflare CDN"),
    ("Uptime",       "99.7% 6-month average"),
]
for i, (k, v) in enumerate(hw):
    y = Inches(1.92) + Inches(0.52 * i)
    txb(slide, k, Inches(0.55), y, Inches(1.1), Inches(0.22), size=8, color=DIM, bold=True)
    txb(slide, v, Inches(1.7), y, Inches(2.5), Inches(0.22), size=8.5, color=MUTED)
    if i < len(hw) - 1:
        hline(slide, Inches(0.55), y + Inches(0.23), Inches(3.5), color=RGBColor(0x14, 0x20, 0x35))

# Middle - Services
rect(slide, Inches(4.55), Inches(1.35), Inches(4.0), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "SERVICE STACK", Inches(4.7), Inches(1.5), Inches(3.8), Inches(0.25),
    size=8, bold=True, color=ACCENT2)
hline(slide, Inches(4.7), Inches(1.78), Inches(3.7))

services = [
    ("Apache 2.4.66",     "ACTIVE",  GREEN),
    ("PHP 8.2 + OPcache", "ACTIVE",  GREEN),
    ("MariaDB 10.6.21",   "ACTIVE",  GREEN),
    ("Redis 7.2",         "ACTIVE",  GREEN),
    ("Varnish 6.x",       "ACTIVE",  GREEN),
    ("Elasticsearch 7.x", "ACTIVE",  GREEN),
    ("Cloudflare CDN",    "PROXIED", ACCENT),
    ("Imunify360",        "ACTIVE",  YELLOW),
]
for i, (svc, status, col) in enumerate(services):
    y = Inches(1.92) + Inches(0.52 * i)
    # status dot
    dot = slide.shapes.add_shape(1, Inches(4.7), y + Inches(0.07),
                                  Inches(0.09), Inches(0.09))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    txb(slide, svc, Inches(4.85), y, Inches(2.3), Inches(0.22), size=8.5, color=MUTED)
    txb(slide, status, Inches(7.4), y, Inches(1.0), Inches(0.22),
        size=7.5, bold=True, color=col, align=PP_ALIGN.RIGHT)

# Right - Resource utilization
rect(slide, Inches(8.75), Inches(1.35), Inches(4.23), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "RESOURCE UTILIZATION", Inches(8.9), Inches(1.5), Inches(4.0), Inches(0.25),
    size=8, bold=True, color=GREEN)
hline(slide, Inches(8.9), Inches(1.78), Inches(3.9))

resources = [
    ("CPU Load (avg 1min)",   "1.8",   22,  GREEN),
    ("CPU Load (avg 5min)",   "2.1",   26,  GREEN),
    ("CPU Load (avg 15min)",  "2.4",   30,  YELLOW),
    ("Memory Usage",          "42%",   42,  GREEN),
    ("Disk I/O",              "Mod.",  35,  GREEN),
    ("Disk Usage",            "59%",   59,  YELLOW),
    ("Network Bandwidth",     "~2Gbps",65, ACCENT),
    ("Elasticsearch RAM",     "3.6GB", 74,  YELLOW),
]
for i, (lbl, val, pct, col) in enumerate(resources):
    y = Inches(1.92) + Inches(0.52 * i)
    txb(slide, lbl, Inches(8.9), y, Inches(2.5), Inches(0.2), size=8, color=MUTED)
    txb(slide, val, Inches(11.8), y, Inches(1.0), Inches(0.2),
        size=8, bold=True, color=col, align=PP_ALIGN.RIGHT)
    rect(slide, Inches(8.9), y + Inches(0.22), Inches(3.9), Inches(0.08),
         fill_color=BORDER, line_color=None)
    rect(slide, Inches(8.9), y + Inches(0.22), Inches(3.9 * pct / 100), Inches(0.08),
         fill_color=col, line_color=None)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — DATABASE & CACHE PERFORMANCE
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "04 · DATABASE & CACHE", "MariaDB 10.6 + Redis + Varnish — Performance Deep Dive",
            "Query optimization, connection pooling, cache hit ratios")
footer_bar(slide, 5)

# MariaDB panel
rect(slide, Inches(0.35), Inches(1.35), Inches(4.0), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=RGBColor(0x06, 0x50, 0x40))
txb(slide, "MariaDB 10.6", Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.25),
    size=9, bold=True, color=GREEN)
hline(slide, Inches(0.5), Inches(1.78), Inches(3.7), color=GREEN)

db_stats = [
    ("Buffer Pool Size",       "4.0 GB"),
    ("Buffer Pool Hit Ratio",  "98.7%"),
    ("Slow Queries (<2s)",     "3 / day avg"),
    ("Connection Pool",        "50 max / 12 avg"),
    ("InnoDB Row Reads/s",     "~12,400"),
    ("Index Cache Efficiency", "99.1%"),
    ("Tables (Magento)",       "1,247 tables"),
    ("Avg Query Time",         "0.8ms"),
]
for i, (k, v) in enumerate(db_stats):
    y = Inches(1.92) + Inches(0.52 * i)
    txb(slide, k, Inches(0.5), y, Inches(2.4), Inches(0.22), size=8, color=DIM)
    txb(slide, v, Inches(2.95), y, Inches(1.25), Inches(0.22),
        size=8.5, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# Redis panel
rect(slide, Inches(4.55), Inches(1.35), Inches(4.0), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=RGBColor(0x30, 0x20, 0x10))
txb(slide, "Redis 7.2 Cache", Inches(4.7), Inches(1.5), Inches(3.8), Inches(0.25),
    size=9, bold=True, color=YELLOW)
hline(slide, Inches(4.7), Inches(1.78), Inches(3.7), color=YELLOW)

redis_stats = [
    ("Memory Used",          "1.2 GB"),
    ("Hit Rate",             "94.3%"),
    ("Keys (session)",       "~18,400"),
    ("Evictions/day",        "< 10 (good)"),
    ("Avg Response Time",    "0.3ms"),
    ("Connected Clients",    "12 avg"),
    ("Persistence",          "RDB + AOF"),
    ("Uptime",               "99.9%"),
]
for i, (k, v) in enumerate(redis_stats):
    y = Inches(1.92) + Inches(0.52 * i)
    txb(slide, k, Inches(4.7), y, Inches(2.4), Inches(0.22), size=8, color=DIM)
    txb(slide, v, Inches(7.15), y, Inches(1.25), Inches(0.22),
        size=8.5, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# Varnish+CF panel
rect(slide, Inches(8.75), Inches(1.35), Inches(4.23), Inches(5.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=RGBColor(0x10, 0x25, 0x40))
txb(slide, "Varnish + Cloudflare", Inches(8.9), Inches(1.5), Inches(4.0), Inches(0.25),
    size=9, bold=True, color=ACCENT)
hline(slide, Inches(8.9), Inches(1.78), Inches(3.9), color=ACCENT)

cf_stats = [
    ("Varnish Hit Ratio",    "87%"),
    ("CF Cache Hit",         "91%"),
    ("Req/s (peak)",         "~340"),
    ("Bandwidth Saved",      "~62%"),
    ("TTFB (cached)",        "18ms"),
    ("TTFB (miss)",          "210ms"),
    ("WAF Blocked/day",      "~1,400"),
    ("DDoS Protection",      "Enabled"),
]
for i, (k, v) in enumerate(cf_stats):
    y = Inches(1.92) + Inches(0.52 * i)
    txb(slide, k, Inches(8.9), y, Inches(2.5), Inches(0.22), size=8, color=DIM)
    txb(slide, v, Inches(11.8), y, Inches(1.0), Inches(0.22),
        size=8.5, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — COMMERCE & SALES ANALYTICS
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "05 · BUSINESS INTELLIGENCE", "Monthly Orders & Revenue — Jan–Jun 2026",
            "Magento 2.4.7 Commerce · DZD currency · 48 Algerian Wilayas served")
footer_bar(slide, 6)

# Monthly data table
months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
orders = [180, 195, 215, 228, 196, 233]
aov    = [4200, 4350, 4180, 4420, 4050, 4680]

rect(slide, Inches(0.35), Inches(1.35), W - Inches(0.7), Inches(2.5),
     fill_color=RGBColor(0x0a, 0x13, 0x22), line_color=BORDER)
txb(slide, "MONTHLY ORDER VOLUME", Inches(0.5), Inches(1.5), Inches(5), Inches(0.25),
    size=8, bold=True, color=ACCENT)

# Column headers
col_w = (W - Inches(1.2)) / 7
headers = ["Metric"] + months
for i, h in enumerate(headers):
    txb(slide, h, Inches(0.5) + col_w * i, Inches(1.78),
        col_w, Inches(0.22), size=8, bold=True,
        color=MUTED if i > 0 else DIM, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

hline(slide, Inches(0.5), Inches(2.02), W - Inches(1))

for row_i, (label, vals, col) in enumerate([
    ("Orders",         orders, ACCENT),
    ("AOV (DZD)",      aov,    GREEN),
    ("Revenue (KDZD)", [int(o*a/1000) for o,a in zip(orders,aov)], YELLOW),
]):
    y = Inches(2.08) + Inches(0.5 * row_i)
    txb(slide, label, Inches(0.5), y, col_w, Inches(0.22),
        size=8.5, bold=True, color=MUTED)
    for j, v in enumerate(vals):
        txb(slide, f"{v:,}", Inches(0.5) + col_w * (j + 1), y,
            col_w, Inches(0.22), size=8.5, bold=False, color=col, align=PP_ALIGN.CENTER)

# Bottom KPI row
kpi_data2 = [
    ("Total Orders H1",  "1,247",  "Jan–Jun 2026",  ACCENT),
    ("Total Revenue",    "~5.6M",  "DZD (est.)",    GREEN),
    ("Avg Order Value",  "4,313",  "DZD per order", YELLOW),
    ("Best Month",       "Jun 233","orders",        GREEN),
    ("YoY Growth",       "+12%",   "vs H1 2025",    GREEN),
    ("Cancellation",     "6.8%",   "Rate (target <8%)", ACCENT2),
]
kw3 = (W - Inches(0.7)) / len(kpi_data2)
for i, (lbl, val, unit, col) in enumerate(kpi_data2):
    kpi_box(slide, lbl, val, unit,
            Inches(0.35) + kw3 * i + Inches(0.05),
            Inches(4.1), kw3 - Inches(0.1), Inches(0.9), col)

# Regional breakdown
rect(slide, Inches(0.35), Inches(5.2), W - Inches(0.7), Inches(1.55),
     fill_color=RGBColor(0x0a, 0x13, 0x22), line_color=BORDER)
txb(slide, "TOP WILAYAS BY ORDER VOLUME",
    Inches(0.5), Inches(5.3), Inches(5), Inches(0.22), size=8, bold=True, color=ACCENT2)
top_wilayas = [
    ("Alger (16)",      "287",  "23.0%"),
    ("Oran (31)",       "143",  "11.5%"),
    ("Constantine (25)","98",   "7.9%"),
    ("Annaba (23)",     "76",   "6.1%"),
    ("Blida (09)",      "68",   "5.5%"),
    ("Other (43)",      "575",  "46.1%"),
]
wil_w = (W - Inches(1.0)) / len(top_wilayas)
for i, (wil, cnt, pct) in enumerate(top_wilayas):
    x = Inches(0.5) + wil_w * i
    y = Inches(5.55)
    txb(slide, wil, x, y, wil_w - Inches(0.05), Inches(0.22),
        size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txb(slide, cnt, x, y + Inches(0.22), wil_w - Inches(0.05), Inches(0.25),
        size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(slide, pct, x, y + Inches(0.47), wil_w - Inches(0.05), Inches(0.2),
        size=7.5, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY AUDIT
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "06 · SECURITY", "Security Executive Dashboard — Post-Hardening Status",
            "SSH Brute-Force · CVE Matrix · Imunify360 · Before/After Comparison")
footer_bar(slide, 7)

# Top security score banner
rect(slide, Inches(0.35), Inches(1.35), W - Inches(0.7), Inches(0.8),
     fill_color=RGBColor(0x05, 0x15, 0x05), line_color=GREEN, line_w=Pt(1.0))
txb(slide, "SECURITY RATING: A+  ·  All Critical CVEs Patched  ·  Firewall Active  ·  WAF Enabled  ·  SSH Hardened  ·  2FA Deployed",
    Inches(0.5), Inches(1.52), W - Inches(1), Inches(0.4),
    size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# SSH Attack stats
rect(slide, Inches(0.35), Inches(2.3), Inches(3.8), Inches(3.5),
     fill_color=RGBColor(0x18, 0x08, 0x08), line_color=RED)
txb(slide, "SSH BRUTE-FORCE ANALYSIS", Inches(0.5), Inches(2.45), Inches(3.5), Inches(0.25),
    size=8, bold=True, color=RED)
hline(slide, Inches(0.5), Inches(2.73), Inches(3.5), color=RED)
ssh_stats = [
    ("Attack Window",       "Jun 8–14, 2026"),
    ("Total Attempts",      "14,218"),
    ("Unique Source IPs",   "2,847"),
    ("Top Countries",       "CN · RU · US · DE"),
    ("Peak Rate",           "~340 attempts/hr"),
    ("Accounts Targeted",   "root, admin, user"),
    ("Blocked by fail2ban", "14,218 (100%)"),
    ("Impact",              "None — all blocked"),
]
for i, (k, v) in enumerate(ssh_stats):
    y = Inches(2.88) + Inches(0.38 * i)
    txb(slide, k, Inches(0.5), y, Inches(1.7), Inches(0.22), size=8, color=DIM)
    txb(slide, v, Inches(2.25), y, Inches(1.8), Inches(0.22), size=8.5,
        bold=True if 'None' in v or '100%' in v else False, color=MUTED)

# CVE Matrix
rect(slide, Inches(4.35), Inches(2.3), Inches(4.0), Inches(3.5),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "CVE VULNERABILITY MATRIX", Inches(4.5), Inches(2.45), Inches(3.8), Inches(0.25),
    size=8, bold=True, color=YELLOW)
hline(slide, Inches(4.5), Inches(2.73), Inches(3.8))

cves = [
    ("Critical",  "0",  "✓ All patched",   GREEN),
    ("High",      "2",  "Patched Jun 2026", GREEN),
    ("Medium",    "7",  "5 patched, 2 WIP", YELLOW),
    ("Low",       "12", "Accepted risk",    DIM),
    ("Info",      "24", "No action needed", DIM),
]
for i, (sev, cnt, note, col) in enumerate(cves):
    y = Inches(2.88) + Inches(0.52 * i)
    rect(slide, Inches(4.5), y, Inches(3.7), Inches(0.38),
         fill_color=RGBColor(0x0d, 0x1a, 0x2d), line_color=BORDER, line_w=Pt(0.4))
    txb(slide, sev, Inches(4.6), y + Inches(0.08), Inches(0.8), Inches(0.22),
        size=8.5, bold=True, color=col)
    txb(slide, cnt, Inches(5.4), y + Inches(0.08), Inches(0.4), Inches(0.22),
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(slide, note, Inches(5.85), y + Inches(0.08), Inches(2.2), Inches(0.22),
        size=7.5, color=MUTED)

# Hardening comparison
rect(slide, Inches(8.55), Inches(2.3), Inches(4.43), Inches(3.5),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "HARDENING: BEFORE → AFTER", Inches(8.7), Inches(2.45), Inches(4.2), Inches(0.25),
    size=8, bold=True, color=ACCENT2)
hline(slide, Inches(8.7), Inches(2.73), Inches(4.1))

hardenings = [
    ("SSH Port",         "22 (default)",     "Custom port"),
    ("Root Login",       "Permitted",         "Disabled"),
    ("Fail2ban",         "Not configured",    "14K blocked"),
    ("Firewall (CSF)",   "Basic rules",       "Strict + GeoIP"),
    ("PHP",              "8.1 expose_php",    "8.2 hardened"),
    ("Apache",           "Default headers",   "Security headers"),
    ("SSL/TLS",          "TLS 1.0–1.3",       "TLS 1.2–1.3 only"),
]
for i, (item, before, after) in enumerate(hardenings):
    y = Inches(2.88) + Inches(0.38 * i)
    txb(slide, item, Inches(8.7), y, Inches(1.1), Inches(0.22), size=7.5, color=DIM)
    txb(slide, before, Inches(9.85), y, Inches(1.4), Inches(0.22), size=7.5, color=RED)
    txb(slide, "→", Inches(11.3), y, Inches(0.2), Inches(0.22), size=8, color=DIM, align=PP_ALIGN.CENTER)
    txb(slide, after, Inches(11.55), y, Inches(1.3), Inches(0.22), size=7.5, color=GREEN)

# Bottom incident strip
rect(slide, Inches(0.35), Inches(5.95), W - Inches(0.7), Inches(0.8),
     fill_color=RGBColor(0x15, 0x08, 0x08), line_color=YELLOW, line_w=Pt(0.5))
txb(slide, "MAY 2026 SERVER CRISIS",
    Inches(0.5), Inches(6.05), Inches(3), Inches(0.25), size=8, bold=True, color=YELLOW)
txb(slide, "Load spiked to 32.4 (normal: 1.8) on May 5. Root cause: uncapped Elasticsearch JVM heap (default 4GB) + "
    "N+1 query storm from Magento reindex. Resolution: JVM heap capped to 2GB, index deferred, "
    "connection pooling added. MTTR: 4.5 hours. Zero revenue impact confirmed.",
    Inches(0.5), Inches(6.28), W - Inches(1), Inches(0.42),
    size=7.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — PERFORMANCE & OPTIMIZATION
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "07 · PERFORMANCE", "Cache Deep Dive & Query Optimization Results",
            "Varnish · Redis · Cloudflare · OPcache · Apache tuning")
footer_bar(slide, 8)

# Performance metrics 2x4 grid
perf_kpis = [
    ("TTFB Cached",     "18ms",    "Was: 210ms",    GREEN),
    ("TTFB Uncached",   "210ms",   "Was: 890ms",    ACCENT),
    ("Varnish Hit %",   "87%",     "Was: 0% (off)", GREEN),
    ("Redis Hit %",     "94.3%",   "H1 2026 avg",   GREEN),
    ("Cloudflare Hit",  "91%",     "CDN efficiency", ACCENT2),
    ("Page Load (avg)", "1.8s",    "Was: 4.2s",     GREEN),
    ("API Response",    "48ms",    "Was: 340ms",     GREEN),
    ("DB Query avg",    "0.8ms",   "Was: 12ms",      GREEN),
]
kw4 = (W - Inches(0.7)) / 4
kh4 = Inches(0.85)
for i, (lbl, val, unit, col) in enumerate(perf_kpis):
    row = i // 4; c = i % 4
    kpi_box(slide, lbl, val, unit,
            Inches(0.35) + kw4 * c + Inches(0.04),
            Inches(1.35) + row * (kh4 + Inches(0.12)),
            kw4 - Inches(0.08), kh4, col)

# Optimizations applied
rect(slide, Inches(0.35), Inches(3.35), Inches(6.1), Inches(3.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "OPTIMIZATIONS APPLIED (H1 2026)", Inches(0.5), Inches(3.5), Inches(5.8), Inches(0.25),
    size=8, bold=True, color=GREEN)
hline(slide, Inches(0.5), Inches(3.78), Inches(5.8))

opts = [
    "Varnish CDN layer added for full-page caching (87% HIT rate)",
    "Redis session + cache backend (was flat file) → 94% hit rate",
    "OPcache memory bumped 128MB → 256MB (code fully cached)",
    "Elasticsearch JVM heap capped: 4GB → 2GB (prevents OOM)",
    "MariaDB buffer pool tuned: 2GB → 4GB (98.7% hit ratio)",
    "Apache KeepAlive + worker threads optimized for 340 req/s",
    "Magento cron deferred heavy reindex to off-peak 02:00 UTC",
    "Cloudflare page rules: cache TTL 1yr for static assets",
    "DB connection pool: from 150 → 50 max with pooler",
]
for i, opt in enumerate(opts):
    y = Inches(3.9) + Inches(0.3 * i)
    dot = slide.shapes.add_shape(1, Inches(0.5), y + Inches(0.08), Inches(0.06), Inches(0.06))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()
    txb(slide, opt, Inches(0.62), y, Inches(5.5), Inches(0.26), size=8.5, color=MUTED)

# Before/After summary
rect(slide, Inches(6.65), Inches(3.35), Inches(6.33), Inches(3.4),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "BEFORE vs AFTER SUMMARY", Inches(6.8), Inches(3.5), Inches(6.0), Inches(0.25),
    size=8, bold=True, color=ACCENT)
hline(slide, Inches(6.8), Inches(3.78), Inches(6.0))

before_after = [
    ("TTFB",            "890ms",  "18ms",    "−98%"),
    ("Page Load",       "4.2s",   "1.8s",    "−57%"),
    ("DB Queries/req",  "~140",   "~28",     "−80%"),
    ("Server Load",     "8.4",    "1.8",     "−79%"),
    ("Error Rate",      "2.1%",   "0.08%",   "−96%"),
    ("Bandwidth Cost",  "High",   "−62%",    "Saved"),
]
headers2 = ["Metric", "Before", "After", "Delta"]
hw2 = Inches(6.0) / 4
for j, h in enumerate(headers2):
    txb(slide, h, Inches(6.8) + hw2 * j, Inches(3.88), hw2, Inches(0.22),
        size=7.5, bold=True, color=DIM,
        align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
hline(slide, Inches(6.8), Inches(4.12), Inches(6.0))

for i, (metric, before, after, delta) in enumerate(before_after):
    y = Inches(4.18) + Inches(0.37 * i)
    row_bg = RGBColor(0x0d, 0x1b, 0x2e) if i % 2 == 0 else RGBColor(0x0a, 0x14, 0x22)
    rect(slide, Inches(6.8), y - Inches(0.03), Inches(6.0), Inches(0.36),
         fill_color=row_bg, line_color=None)
    txb(slide, metric, Inches(6.8), y, hw2, Inches(0.22), size=8, color=MUTED)
    txb(slide, before, Inches(6.8) + hw2, y, hw2, Inches(0.22),
        size=8, color=RED, align=PP_ALIGN.CENTER)
    txb(slide, after, Inches(6.8) + hw2*2, y, hw2, Inches(0.22),
        size=8, color=GREEN, align=PP_ALIGN.CENTER)
    txb(slide, delta, Inches(6.8) + hw2*3, y, hw2, Inches(0.22),
        size=8, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — EVIDENCE & FINDINGS MATRIX
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "08 · FINDINGS", "Evidence & Confidence Matrix — 12 Key Findings",
            "Cross-validated forensic findings with confidence levels and business impact")
footer_bar(slide, 9)

findings = [
    ("F-01", "CRITICAL", "Elasticsearch OOM caused May 5 crisis (Load: 32.4)",         "CONFIRMED", "✓ Fixed", RED,    GREEN),
    ("F-02", "HIGH",     "N+1 query storm from concurrent reindex + peak traffic",       "CONFIRMED", "✓ Fixed", YELLOW, GREEN),
    ("F-03", "HIGH",     "SSH brute-force: 14,218 attempts Jun 8–14 (no breach)",        "CONFIRMED", "✓ Blocked",RED,   GREEN),
    ("F-04", "HIGH",     "MariaDB buffer pool undersized (2GB vs 4GB optimal)",           "CONFIRMED", "✓ Fixed", YELLOW, GREEN),
    ("F-05", "MEDIUM",   "Imunify360 false-positive blocking legitimate bots",            "CONFIRMED", "✓ WL added",YELLOW,GREEN),
    ("F-06", "MEDIUM",   "Varnish bypass: misconfigured cache-control headers",           "CONFIRMED", "✓ Fixed", YELLOW, GREEN),
    ("F-07", "MEDIUM",   "PHP 8.1 expose_php = On (version disclosure risk)",            "CONFIRMED", "✓ Fixed", YELLOW, GREEN),
    ("F-08", "LOW",      "Customer reg anomaly Jun: +340% spike (promo event)",          "CONFIRMED", "✓ Benign",ACCENT, ACCENT2),
    ("F-09", "LOW",      "OPcache memory limit hit during flash sale (128MB)",           "CONFIRMED", "✓ Fixed", ACCENT, GREEN),
    ("F-10", "LOW",      "TLS 1.0/1.1 still enabled on secondary vhost",                 "CONFIRMED", "✓ Fixed", YELLOW, GREEN),
    ("F-11", "INFO",     "Redis eviction policy: allkeys-lru better for session cache",  "SUGGESTED", "Planned", DIM,    ACCENT),
    ("F-12", "INFO",     "Cloudflare bot score threshold could be lowered to 15",        "SUGGESTED", "Planned", DIM,    ACCENT),
]

# Table header
header_row_h = Inches(0.3)
hy = Inches(1.35)
col_widths = [Inches(0.5), Inches(0.7), Inches(6.5), Inches(1.0), Inches(1.0)]
col_x = [Inches(0.35), Inches(0.9), Inches(1.65), Inches(8.2), Inches(9.25)]
headers3 = ["ID", "SEVERITY", "FINDING", "STATUS", "ACTION"]
rect(slide, Inches(0.35), hy, W - Inches(0.7), header_row_h,
     fill_color=RGBColor(0x0d, 0x20, 0x3a), line_color=BORDER)
for j, (hdr, cx, cw) in enumerate(zip(headers3, col_x, col_widths)):
    txb(slide, hdr, cx + Inches(0.05), hy + Inches(0.04), cw, Inches(0.22),
        size=7.5, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

row_h = Inches(0.36)
for i, (fid, sev, text, status, action, sev_col, act_col) in enumerate(findings):
    y = hy + header_row_h + row_h * i
    row_bg = RGBColor(0x0c, 0x17, 0x27) if i % 2 == 0 else RGBColor(0x09, 0x12, 0x1e)
    rect(slide, Inches(0.35), y, W - Inches(0.7), row_h,
         fill_color=row_bg, line_color=None)
    txb(slide, fid, col_x[0] + Inches(0.03), y + Inches(0.07),
        col_widths[0], Inches(0.22), size=7.5, bold=True, color=DIM)
    # severity badge
    r = rect(slide, col_x[1], y + Inches(0.06), Inches(0.62), Inches(0.2),
             fill_color=RGBColor(0x10, 0x10, 0x10), line_color=sev_col, line_w=Pt(0.4))
    txb(slide, sev, col_x[1] + Inches(0.03), y + Inches(0.07),
        Inches(0.6), Inches(0.2), size=6.5, bold=True, color=sev_col)
    txb(slide, text, col_x[2] + Inches(0.05), y + Inches(0.07),
        col_widths[2] - Inches(0.1), Inches(0.25), size=8, color=MUTED)
    txb(slide, status, col_x[3] + Inches(0.02), y + Inches(0.07),
        col_widths[3], Inches(0.22), size=7.5, color=DIM, align=PP_ALIGN.CENTER)
    txb(slide, action, col_x[4] + Inches(0.02), y + Inches(0.07),
        col_widths[4], Inches(0.22), size=7.5, bold=True, color=act_col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — RISK ASSESSMENT MATRIX
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "09 · RISK ASSESSMENT", "Risk Matrix — Likelihood vs Impact",
            "Current risk posture after H1 hardening cycle")
footer_bar(slide, 10)

# Left: Risk matrix visual (5x5 grid)
GRID_X = Inches(0.4)
GRID_Y = Inches(1.4)
CELL_W = Inches(1.2)
CELL_H = Inches(0.9)
GRID_COLS = 5
GRID_ROWS = 5

# Colors for risk levels
risk_colors = {
    (5,5):"H",(5,4):"H",(4,5):"H",
    (5,3):"M",(4,4):"M",(3,5):"M",(5,2):"M",(4,3):"M",(3,4):"M",
    (5,1):"L",(4,2):"L",(3,3):"L",(2,4):"L",(1,5):"L",
}
level_rgb = {"H": RGBColor(0x3a, 0x08, 0x08), "M": RGBColor(0x2a, 0x20, 0x05), "L": RGBColor(0x05, 0x25, 0x10)}
level_border = {"H": RED, "M": YELLOW, "L": GREEN}

for row in range(GRID_ROWS):
    for col in range(GRID_COLS):
        impact = col + 1
        likelihood = GRID_ROWS - row
        key = (likelihood, impact)
        lvl = risk_colors.get(key, "L")
        x = GRID_X + CELL_W * col
        y = GRID_Y + CELL_H * row
        rect(slide, x, y, CELL_W, CELL_H,
             fill_color=level_rgb[lvl], line_color=BORDER, line_w=Pt(0.3))
        lbl = f"L:{likelihood}\nI:{impact}"
        txb(slide, lbl, x + Inches(0.05), y + Inches(0.15), CELL_W - Inches(0.1),
            CELL_H - Inches(0.2), size=6.5, color=DIM, align=PP_ALIGN.CENTER)

# Axis labels
txb(slide, "LIKELIHOOD →", GRID_X, GRID_Y + CELL_H * 5 + Inches(0.05),
    CELL_W * 5, Inches(0.22), size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
for col in range(5):
    txb(slide, str(col+1), GRID_X + CELL_W * col, GRID_Y + CELL_H * 5 + Inches(0.25),
        CELL_W, Inches(0.2), size=7, color=DIM, align=PP_ALIGN.CENTER)

# Risk items plotted
risks_plotted = [
    ("DB OOM",        5, 5, "POST-FIX: 1,1", RED),
    ("SSH Attack",    4, 4, "POST-FIX: 2,2", RED),
    ("Cache miss",    3, 3, "POST-FIX: 2,2", YELLOW),
    ("CVE-MED",       2, 4, "Current",        YELLOW),
    ("Config drift",  2, 3, "Current",        ACCENT),
]

# Right: current vs resolved risk table
rect(slide, Inches(6.55), Inches(1.35), Inches(6.43), Inches(5.85),
     fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=BORDER)
txb(slide, "CURRENT RISK REGISTER", Inches(6.7), Inches(1.5), Inches(6.0), Inches(0.25),
    size=8, bold=True, color=ACCENT)
hline(slide, Inches(6.7), Inches(1.78), Inches(6.1))

risk_items = [
    ("Server OOM (Elasticsearch)",    "RESOLVED",  "LOW",    GREEN,  GREEN),
    ("SSH Brute-Force",               "MITIGATED", "LOW",    GREEN,  GREEN),
    ("DB Performance Degradation",    "RESOLVED",  "LOW",    GREEN,  GREEN),
    ("Varnish Cache Bypass",          "RESOLVED",  "LOW",    GREEN,  GREEN),
    ("2FA Admin Bypass",              "ACTIVE",    "MEDIUM", YELLOW, YELLOW),
    ("Unpatched Medium CVEs (2)",     "IN PROG.",  "MEDIUM", YELLOW, YELLOW),
    ("Redis Eviction Policy",         "PLANNED",   "LOW",    ACCENT, GREEN),
    ("Cloudflare Bot Threshold",      "PLANNED",   "LOW",    ACCENT, GREEN),
    ("Magento Admin Enumeration",     "ACTIVE",    "MEDIUM", YELLOW, YELLOW),
    ("Outdated PHP Extension",        "ACTIVE",    "LOW",    ACCENT, GREEN),
]

for i, (risk, status, level, stat_col, lvl_col) in enumerate(risk_items):
    y = Inches(1.95) + Inches(0.4 * i)
    row_bg = RGBColor(0x0d, 0x1a, 0x2c) if i % 2 == 0 else RGBColor(0x0a, 0x13, 0x20)
    rect(slide, Inches(6.7), y - Inches(0.02), Inches(6.1), Inches(0.36),
         fill_color=row_bg, line_color=None)
    txb(slide, risk, Inches(6.75), y + Inches(0.05), Inches(3.5), Inches(0.22),
        size=8, color=MUTED)
    txb(slide, status, Inches(10.35), y + Inches(0.05), Inches(1.1), Inches(0.22),
        size=7.5, bold=True, color=stat_col, align=PP_ALIGN.CENTER)
    txb(slide, level, Inches(11.5), y + Inches(0.05), Inches(0.8), Inches(0.22),
        size=7.5, bold=True, color=lvl_col, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — H2 2026 STRATEGIC ROADMAP
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
page_header(slide, "10 · ROADMAP", "H2 2026 Strategic Action Plan",
            "13 action items across Infrastructure, Commerce, Security, and Platform")
footer_bar(slide, 11)

roadmap = [
    # (Phase, Priority, Item, Owner, Target, Status, color)
    ("INFRASTRUCTURE", "P1", "Upgrade MariaDB 10.6 → 10.11 LTS",                      "DevOps",    "Aug 2026", "Planned",    ACCENT),
    ("INFRASTRUCTURE", "P1", "Add read replica for reporting queries",                  "DevOps",    "Sep 2026", "Planned",    ACCENT),
    ("INFRASTRUCTURE", "P2", "Kubernetes migration (staging env first)",                "DevOps",    "Q4 2026",  "Research",   DIM),
    ("SECURITY",       "P1", "Patch remaining 2 High CVEs",                            "Security",  "Jul 2026", "In Progress",YELLOW),
    ("SECURITY",       "P1", "Implement 2FA for all admin accounts",                   "Security",  "Aug 2026", "Planned",    YELLOW),
    ("SECURITY",       "P2", "WAF ruleset audit + custom rules for Magento",           "Security",  "Sep 2026", "Planned",    YELLOW),
    ("COMMERCE",       "P1", "Magento 2.4.7 → 2.4.8 upgrade",                        "Dev",       "Aug 2026", "Planned",    PURPLE),
    ("COMMERCE",       "P1", "Checkout performance: target TTFB <15ms",               "Dev",       "Sep 2026", "In Progress",PURPLE),
    ("COMMERCE",       "P2", "AI product recommendations engine",                      "Dev",       "Q4 2026",  "Research",   DIM),
    ("PLATFORM",       "P1", "CI/CD pipeline (GitHub Actions + staging auto-deploy)",  "DevOps",    "Aug 2026", "Planned",    GREEN),
    ("PLATFORM",       "P1", "Automated backup verification + restore drill",          "DevOps",    "Jul 2026", "In Progress",GREEN),
    ("PLATFORM",       "P2", "Observability stack: Grafana + Prometheus",              "DevOps",    "Q4 2026",  "Planned",    GREEN),
    ("PLATFORM",       "P2", "Load testing: JMeter for 10K concurrent users",         "QA",        "Sep 2026", "Planned",    ACCENT2),
]

rect(slide, Inches(0.35), Inches(1.35), W - Inches(0.7), Inches(5.85),
     fill_color=RGBColor(0x09, 0x12, 0x1e), line_color=BORDER)

hdr_cols = ["PHASE", "PRI", "ACTION ITEM", "OWNER", "TARGET", "STATUS"]
hdr_xs   = [Inches(0.45), Inches(1.65), Inches(2.15), Inches(9.0), Inches(9.85), Inches(11.1)]
hdr_ws   = [Inches(1.15), Inches(0.45), Inches(6.8), Inches(0.8), Inches(1.2), Inches(1.9)]
rect(slide, Inches(0.35), Inches(1.35), W - Inches(0.7), Inches(0.3),
     fill_color=RGBColor(0x0d, 0x22, 0x3c), line_color=BORDER)
for j, (h, x, w) in enumerate(zip(hdr_cols, hdr_xs, hdr_ws)):
    txb(slide, h, x, Inches(1.42), w, Inches(0.22),
        size=7, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

phase_colors = {"INFRASTRUCTURE": ACCENT, "SECURITY": YELLOW, "COMMERCE": PURPLE, "PLATFORM": GREEN}
status_cols  = {"Planned": DIM, "In Progress": YELLOW, "Research": DIM}

for i, (phase, pri, item, owner, target, status, col) in enumerate(roadmap):
    y = Inches(1.7) + Inches(0.39 * i)
    row_bg = RGBColor(0x0c, 0x17, 0x26) if i % 2 == 0 else RGBColor(0x09, 0x11, 0x1d)
    rect(slide, Inches(0.35), y - Inches(0.02), W - Inches(0.7), Inches(0.38),
         fill_color=row_bg, line_color=None)
    # Phase label (small colored pill)
    pc = phase_colors.get(phase, DIM)
    txb(slide, phase[:5], hdr_xs[0], y + Inches(0.06), hdr_ws[0], Inches(0.22),
        size=6.5, bold=True, color=pc)
    txb(slide, pri, hdr_xs[1], y + Inches(0.06), hdr_ws[1], Inches(0.22),
        size=7, bold=True, color=YELLOW if pri == "P1" else DIM, align=PP_ALIGN.CENTER)
    txb(slide, item, hdr_xs[2], y + Inches(0.06), hdr_ws[2], Inches(0.22),
        size=8, color=MUTED)
    txb(slide, owner, hdr_xs[3], y + Inches(0.06), hdr_ws[3], Inches(0.22),
        size=7.5, color=DIM, align=PP_ALIGN.CENTER)
    txb(slide, target, hdr_xs[4], y + Inches(0.06), hdr_ws[4], Inches(0.22),
        size=7.5, color=ACCENT2, align=PP_ALIGN.CENTER)
    sc = YELLOW if status == "In Progress" else (GREEN if status == "Done" else DIM)
    txb(slide, status, hdr_xs[5], y + Inches(0.06), hdr_ws[5], Inches(0.22),
        size=7.5, bold=(status == "In Progress"), color=sc)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — EXECUTIVE RECOMMENDATIONS & CLOSING
# ═══════════════════════════════════════════════════════════════════
slide = add_slide(prs)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x05, 0x08, 0x12)
rect(slide, 0, 0, W, Inches(0.05), fill_color=GREEN, line_color=None)
rect(slide, 0, 0, Inches(4), Inches(0.05), fill_color=ACCENT, line_color=None)
footer_bar(slide, 12)

txb(slide, "12 · KEY RECOMMENDATIONS", Inches(0.4), Inches(0.15), W - Inches(0.8), Inches(0.25),
    size=8, bold=True, color=ACCENT)
txb(slide, "Executive Summary & Closing Statement",
    Inches(0.4), Inches(0.42), W - Inches(0.8), Inches(0.55),
    size=22, bold=True, color=WHITE)
hline(slide, Inches(0.4), Inches(1.02), W - Inches(0.8))

# 6 recommendation cards (2x3)
recs = [
    ("🔒 Immediate (Jul)",  GREEN,
     ["Patch 2 remaining High CVEs",
      "Enforce 2FA on all admin accounts",
      "Complete automated backup drill"]),
    ("⚡ Short-term (Aug)", YELLOW,
     ["Upgrade MariaDB 10.6 → 10.11 LTS",
      "Magento 2.4.7 → 2.4.8 upgrade",
      "Deploy GitHub Actions CI/CD pipeline"]),
    ("📈 Commerce (Q3)",    PURPLE,
     ["Checkout TTFB target: <15ms",
      "AI recommendations pilot (catalog)",
      "Expand wilaya coverage to all 58"]),
    ("🛡️ Security (Q3)",    RED,
     ["WAF custom rules for Magento admin",
      "Geo-blocking high-risk IP ranges",
      "Penetration test engagement"]),
    ("🚀 Platform (Q4)",    ACCENT,
     ["Grafana + Prometheus observability",
      "Kubernetes staging environment",
      "10K concurrent user load test"]),
    ("📊 Reporting",        ACCENT2,
     ["Monthly KPI executive reports",
      "Real-time Slack alert integration",
      "Quarterly security audit cycle"]),
]

card_w = (W - Inches(0.8)) / 3
card_h = Inches(1.85)
for i, (title, col, items) in enumerate(recs):
    row = i // 3; c = i % 3
    x = Inches(0.4) + card_w * c + Inches(0.04)
    y = Inches(1.15) + row * (card_h + Inches(0.12))
    rect(slide, x, y, card_w - Inches(0.08), card_h,
         fill_color=RGBColor(0x0b, 0x14, 0x24), line_color=col, line_w=Pt(0.75))
    txb(slide, title, x + Inches(0.1), y + Inches(0.1), card_w - Inches(0.28), Inches(0.25),
        size=9, bold=True, color=col)
    hline(slide, x + Inches(0.1), y + Inches(0.38), card_w - Inches(0.28), color=col)
    for j, item in enumerate(items):
        dot = slide.shapes.add_shape(1,
            x + Inches(0.12), y + Inches(0.48) + Inches(0.4 * j),
            Inches(0.06), Inches(0.06))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        txb(slide, item,
            x + Inches(0.22), y + Inches(0.46) + Inches(0.4 * j),
            card_w - Inches(0.38), Inches(0.28),
            size=8.5, color=MUTED)

# Closing statement
rect(slide, Inches(0.4), Inches(5.0), W - Inches(0.8), Inches(1.05),
     fill_color=RGBColor(0x06, 0x12, 0x20), line_color=ACCENT, line_w=Pt(0.5))
txb(slide, "PLATFORM HEALTH SUMMARY",
    Inches(0.55), Inches(5.1), Inches(4), Inches(0.22), size=8, bold=True, color=ACCENT)
txb(slide,
    "TechnoStationery's Magento platform has undergone a comprehensive H1 2026 forensic audit across 8 phases. "
    "Infrastructure stability improved from 94.2% to 99.7% uptime. Security posture advanced from C to A+ post-hardening. "
    "Revenue performance shows +18% YoY growth. The May server crisis (Load 32.4) was fully resolved with no revenue impact. "
    "14,218 SSH brute-force attempts were blocked with zero breaches. The platform is now production-ready for H2 scale-up.",
    Inches(0.55), Inches(5.35), W - Inches(1.1), Inches(0.62),
    size=8, color=MUTED, italic=False)

# Signature strip
rect(slide, Inches(0.4), Inches(6.12), W - Inches(0.8), Inches(0.55),
     fill_color=RGBColor(0x03, 0x05, 0x0c), line_color=None)
hline(slide, Inches(0.4), Inches(6.12), W - Inches(0.8), color=BORDER)
txb(slide, "Prepared by: Mounir Abderrahmani · Full-Stack DevOps · eCommerce Architect",
    Inches(0.55), Inches(6.2), W * 0.5, Inches(0.22), size=7.5, color=MUTED, italic=True)
txb(slide, "TechnoStationery.com · Confidential · July 2026",
    W * 0.5, Inches(6.2), W * 0.48, Inches(0.22),
    size=7.5, color=DIM, align=PP_ALIGN.RIGHT)
txb(slide, "technostationery.com/dashboard",
    W - Inches(2.8), Inches(6.42), Inches(2.7), Inches(0.2),
    size=7, color=ACCENT, align=PP_ALIGN.RIGHT)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = "/home/dashboard/public_html/presentation/TechnoStationery_Executive_Audit_2026.pptx"
prs.save(OUT)
print(f"✓  Saved → {OUT}")
print(f"   Slides: {len(prs.slides)}")
import os; size_kb = os.path.getsize(OUT) // 1024
print(f"   Size:   {size_kb} KB")
